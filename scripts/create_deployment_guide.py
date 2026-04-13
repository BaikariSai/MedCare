#!/usr/bin/env python3
"""
MedCare Hospital - Deployment Guide PPT Generator
Creates a comprehensive deployment guide for Local + Vercel + Supabase
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============== COLORS ==============
BLUE_DARK = RGBColor(0x0F, 0x17, 0x2A)
BLUE_PRIMARY = RGBColor(0x1D, 0x4E, 0xD8)
BLUE_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)
BLUE_ACCENT = RGBColor(0x3B, 0x82, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_900 = RGBColor(0x11, 0x18, 0x27)
GRAY_700 = RGBColor(0x37, 0x41, 0x51)
GRAY_500 = RGBColor(0x6B, 0x72, 0x80)
GRAY_200 = RGBColor(0xE5, 0xE7, 0xEB)
GRAY_100 = RGBColor(0xF3, 0xF4, 0xF6)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
TEAL = RGBColor(0x0D, 0x94, 0x88)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=GRAY_900, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=GRAY_700, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_code_block(slide, left, top, width, height, code, font_size=11):
    shape = add_shape(slide, left, top, width, height, RGBColor(0x1E, 0x1E, 0x2E))
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0xA6, 0xE2, 0x2E)
        p.font.name = "Consolas"
        p.space_after = Pt(2)
    return txBox


def add_error_card(slide, left, top, width, error_title, error_msg, solution, idx=0):
    card_h = Inches(1.6)
    # Red accent bar
    add_shape(slide, left, top, Inches(0.08), card_h, RED)
    # Card bg
    add_shape(slide, left + Inches(0.08), top, width - Inches(0.08), card_h, WHITE, GRAY_200)
    # Error title
    add_text_box(slide, left + Inches(0.3), top + Inches(0.1), width - Inches(0.5), Inches(0.35),
                 f"Error #{idx+1}: {error_title}", font_size=13, color=RED, bold=True)
    # Error message
    add_code_block(slide, left + Inches(0.3), top + Inches(0.45), width - Inches(0.5), Inches(0.4), error_msg, font_size=9)
    # Solution
    add_text_box(slide, left + Inches(0.3), top + Inches(0.9), width - Inches(0.5), Inches(0.6),
                 f"Fix: {solution}", font_size=11, color=GREEN)
    return card_h


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)

# Decorative shapes
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), BLUE_PRIMARY)
add_shape(slide, Inches(-1), Inches(5.5), Inches(6), Inches(6), RGBColor(0x1D, 0x2D, 0x50))
add_shape(slide, Inches(9), Inches(-1), Inches(6), Inches(6), RGBColor(0x1D, 0x2D, 0x50))

# Heart icon (simulated with shape)
shape = add_shape(slide, Inches(6.15), Inches(1.2), Inches(1), Inches(1), BLUE_PRIMARY)
add_text_box(slide, Inches(6.15), Inches(1.3), Inches(1), Inches(0.8), "+", font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.2),
             "MedCare Hospital", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10.3), Inches(0.6),
             "Deployment Guide", font_size=32, color=BLUE_ACCENT, bold=False, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.5),
             "Local Development  |  Vercel  |  Supabase", font_size=18, color=GRAY_500, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.4),
             "Full-Stack Next.js + Supabase (PostgreSQL) Hospital Management System", font_size=14, color=GRAY_500, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Table of Contents
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE_PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "Table of Contents", font_size=32, color=BLUE_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.15), Inches(2), Inches(0.05), BLUE_PRIMARY)

toc_items = [
    ("01", "Architecture Overview", "System architecture & tech stack"),
    ("02", "Prerequisites", "Required tools, accounts & credentials"),
    ("03", "Supabase Database Setup", "Create tables, configure RLS, seed data"),
    ("04", "Local Deployment", "Step-by-step local development setup"),
    ("05", "Vercel Deployment", "Production deployment to Vercel"),
    ("06", "Environment Variables", "Complete env config reference"),
    ("07", "Common Errors & Fixes", "Troubleshooting deployment issues"),
    ("08", "Post-Deployment Checklist", "Verification & monitoring"),
]

for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(1.6) + Inches(i * 0.68)
    add_shape(slide, Inches(0.8), y, Inches(0.7), Inches(0.5), BLUE_PRIMARY if i < 3 else BLUE_LIGHT)
    add_text_box(slide, Inches(0.8), y + Inches(0.07), Inches(0.7), Inches(0.4),
                 num, font_size=18, color=WHITE if i < 3 else BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.7), y + Inches(0.02), Inches(5), Inches(0.3),
                 title, font_size=16, color=GRAY_900, bold=True)
    add_text_box(slide, Inches(1.7), y + Inches(0.28), Inches(5), Inches(0.25),
                 desc, font_size=11, color=GRAY_500)

# ============================================================
# SLIDE 3: Architecture Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE_PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "01  Architecture Overview", font_size=28, color=BLUE_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05), BLUE_PRIMARY)

# Architecture boxes
boxes = [
    ("Frontend", "Next.js 14 (App Router)\nTailwind CSS\nShadCN UI\nLucide Icons", BLUE_PRIMARY, Inches(0.8)),
    ("Backend", "Next.js API Routes\nCatch-all Route Handler\nServer-side Logic\nJWT Auth", TEAL, Inches(3.5)),
    ("Database", "Supabase (PostgreSQL)\n6 Tables\nRow Level Security\nReal-time capable", GREEN, Inches(6.2)),
    ("Hosting", "Vercel (Frontend+API)\nAutomatic SSL\nEdge Network\nCI/CD Pipeline", ORANGE, Inches(8.9)),
]
for label, desc, color, left in boxes:
    add_shape(slide, left, Inches(1.6), Inches(2.5), Inches(2.6), color)
    add_text_box(slide, left + Inches(0.2), Inches(1.8), Inches(2.1), Inches(0.4),
                 label, font_size=18, color=WHITE, bold=True)
    add_shape(slide, left + Inches(0.2), Inches(2.25), Inches(2.1), Inches(0.03), WHITE)
    add_text_box(slide, left + Inches(0.2), Inches(2.4), Inches(2.1), Inches(1.5),
                 desc, font_size=12, color=WHITE)

# DB Tables
add_text_box(slide, Inches(0.8), Inches(4.6), Inches(11), Inches(0.5),
             "Database Tables (PostgreSQL via Supabase)", font_size=16, color=GRAY_900, bold=True)
tables = ["departments", "doctors", "appointments", "users", "blogs", "contact_queries"]
for i, t in enumerate(tables):
    x = Inches(0.8) + Inches(i * 2.05)
    add_shape(slide, x, Inches(5.1), Inches(1.85), Inches(0.45), BLUE_LIGHT)
    add_text_box(slide, x, Inches(5.15), Inches(1.85), Inches(0.35),
                 t, font_size=12, color=BLUE_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)

# Tech stack
add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.4),
             "Tech Stack:  Next.js 14  •  React 18  •  Tailwind CSS  •  ShadCN UI  •  Supabase JS v2  •  PostgreSQL  •  Sonner  •  Lucide React  •  UUID  •  Zod",
             font_size=11, color=GRAY_500, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: Prerequisites
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE_PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "02  Prerequisites", font_size=28, color=BLUE_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05), BLUE_PRIMARY)

# Left column - Tools
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.2), GRAY_100)
add_text_box(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(0.4),
             "Required Tools", font_size=18, color=GRAY_900, bold=True)
tools = [
    "Node.js v18+ (LTS recommended)",
    "Yarn package manager (npm install -g yarn)",
    "Git for version control",
    "Code editor (VS Code recommended)",
    "Terminal / Command Line",
]
add_bullet_list(slide, Inches(1.1), Inches(2.2), Inches(5.2), Inches(2.5), 
                [f"  {t}" for t in tools], font_size=13, color=GRAY_700)

add_text_box(slide, Inches(1.1), Inches(4.3), Inches(5), Inches(0.4),
             "Verify Installation:", font_size=14, color=GRAY_900, bold=True)
add_code_block(slide, Inches(1.1), Inches(4.7), Inches(5.2), Inches(1.5),
               "node --version    # v18.x or higher\nyarn --version    # 1.22.x\ngit --version     # 2.x\nnpx --version     # 10.x")

# Right column - Accounts
add_shape(slide, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.2), GRAY_100)
add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5), Inches(0.4),
             "Required Accounts & Credentials", font_size=18, color=GRAY_900, bold=True)
accounts = [
    "Supabase Account (free tier available)",
    "  supabase.com/dashboard",
    "  You need: Project URL, Anon Key, Service Role Key",
    "",
    "Vercel Account (for production deployment)",
    "  vercel.com",
    "  Connect your GitHub repository",
    "",
    "GitHub Account (optional but recommended)",
    "  For CI/CD with Vercel",
]
add_bullet_list(slide, Inches(7.2), Inches(2.2), Inches(5), Inches(4), 
                accounts, font_size=12, color=GRAY_700, spacing=Pt(4))

# ============================================================
# SLIDE 5: Supabase Database Setup
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE_PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "03  Supabase Database Setup", font_size=28, color=BLUE_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05), BLUE_PRIMARY)

# Step 1
add_shape(slide, Inches(0.8), Inches(1.5), Inches(0.5), Inches(0.5), BLUE_PRIMARY)
add_text_box(slide, Inches(0.8), Inches(1.55), Inches(0.5), Inches(0.4), "1", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(1.55), Inches(5), Inches(0.4),
             "Create a Supabase Project", font_size=16, color=GRAY_900, bold=True)
add_text_box(slide, Inches(1.5), Inches(1.9), Inches(5), Inches(0.6),
             "Go to supabase.com/dashboard > New Project > Choose region > Set database password", font_size=12, color=GRAY_500)

# Step 2
add_shape(slide, Inches(0.8), Inches(2.5), Inches(0.5), Inches(0.5), BLUE_PRIMARY)
add_text_box(slide, Inches(0.8), Inches(2.55), Inches(0.5), Inches(0.4), "2", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(2.55), Inches(5), Inches(0.4),
             "Get Your API Credentials", font_size=16, color=GRAY_900, bold=True)
add_text_box(slide, Inches(1.5), Inches(2.9), Inches(5), Inches(0.4),
             "Settings > API > Copy: Project URL, anon key, service_role key", font_size=12, color=GRAY_500)

# Step 3
add_shape(slide, Inches(0.8), Inches(3.4), Inches(0.5), Inches(0.5), BLUE_PRIMARY)
add_text_box(slide, Inches(0.8), Inches(3.45), Inches(0.5), Inches(0.4), "3", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.45), Inches(5), Inches(0.4),
             "Create Tables (SQL Editor)", font_size=16, color=GRAY_900, bold=True)

# SQL Code block
sql_code = """CREATE TABLE IF NOT EXISTS departments (
  id TEXT PRIMARY KEY, name TEXT NOT NULL,
  slug TEXT NOT NULL, icon TEXT,
  description TEXT, services JSONB DEFAULT '[]'::jsonb
);
CREATE TABLE IF NOT EXISTS doctors (
  id TEXT PRIMARY KEY, name TEXT NOT NULL,
  slug TEXT, department_id TEXT, department TEXT,
  qualification TEXT, experience INTEGER DEFAULT 0,
  bio TEXT, gender TEXT, rating DECIMAL(3,1),
  reviews INTEGER DEFAULT 0, featured BOOLEAN DEFAULT false,
  availability JSONB, consultation_fee INTEGER DEFAULT 0
);
-- appointments, users, blogs, contact_queries
-- (See full SQL in project /api/setup-sql endpoint)"""

add_code_block(slide, Inches(0.8), Inches(4.0), Inches(5.8), Inches(3.0), sql_code, font_size=9)

# Right side - credentials
add_shape(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.5), RGBColor(0xFE, 0xF3, 0xC7))
add_text_box(slide, Inches(7.3), Inches(1.7), Inches(5), Inches(0.4),
             "Important: Save These Credentials", font_size=16, color=ORANGE, bold=True)
creds = [
    "SUPABASE_URL",
    "  https://[project-ref].supabase.co",
    "  Found in: Settings > API > URL",
    "",
    "SUPABASE_ANON_KEY",
    "  eyJhbGciOi... (public key)",
    "  Found in: Settings > API > anon public",
    "",
    "SUPABASE_SERVICE_ROLE_KEY",
    "  eyJhbGciOi... (secret key)",
    "  Found in: Settings > API > service_role",
    "  NEVER expose in client-side code!",
    "",
    "After running SQL, call /api/seed to",
    "populate initial data automatically.",
]
add_bullet_list(slide, Inches(7.3), Inches(2.2), Inches(4.8), Inches(4.5), creds, font_size=11, color=GRAY_700, spacing=Pt(3))

# ============================================================
# SLIDE 6: Local Deployment
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE_PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "04  Local Deployment (Step-by-Step)", font_size=28, color=BLUE_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05), BLUE_PRIMARY)

# Left column - Steps
steps_code = """# Step 1: Clone the repository
git clone https://github.com/your-org/medcare-hospital.git
cd medcare-hospital

# Step 2: Install dependencies
yarn install

# Step 3: Create .env file
cp .env.example .env

# Step 4: Edit .env with your Supabase credentials
SUPABASE_URL=https://xxx.supabase.co
SUPABASE_ANON_KEY=eyJhbGciOi...
SUPABASE_SERVICE_ROLE_KEY=eyJhbGciOi...
NEXT_PUBLIC_BASE_URL=http://localhost:3000

# Step 5: Run development server
yarn dev

# Step 6: Open browser
# http://localhost:3000
# First visit triggers /api/seed automatically"""

add_code_block(slide, Inches(0.8), Inches(1.4), Inches(6.5), Inches(5.3), steps_code, font_size=11)

# Right column - Notes
add_shape(slide, Inches(7.6), Inches(1.4), Inches(5), Inches(2.2), BLUE_LIGHT)
add_text_box(slide, Inches(7.9), Inches(1.6), Inches(4.5), Inches(0.4),
             "Local Dev Notes", font_size=16, color=BLUE_PRIMARY, bold=True)
notes = [
    "Hot reload is enabled by default",
    "API routes: http://localhost:3000/api/*",
    "Frontend: http://localhost:3000",
    "Use yarn (NOT npm) for consistency",
    "Server restarts needed only for .env changes",
]
add_bullet_list(slide, Inches(7.9), Inches(2.1), Inches(4.5), Inches(1.4), 
                [f"  {n}" for n in notes], font_size=11, color=GRAY_700, spacing=Pt(4))

add_shape(slide, Inches(7.6), Inches(3.8), Inches(5), Inches(2.8), RGBColor(0xDC, 0xFC, 0xE7))
add_text_box(slide, Inches(7.9), Inches(4.0), Inches(4.5), Inches(0.4),
             "Project Structure", font_size=16, color=GREEN, bold=True)
structure = """/app
 ├── app/
 │   ├── api/[[...path]]/route.js  # Backend
 │   ├── page.js                    # Frontend
 │   ├── layout.js                  # Layout + SEO
 │   └── globals.css                # Styles
 ├── components/ui/                 # ShadCN
 ├── .env                           # Credentials
 ├── package.json
 └── tailwind.config.js"""
add_code_block(slide, Inches(7.9), Inches(4.5), Inches(4.5), Inches(2.0), structure, font_size=10)

# ============================================================
# SLIDE 7: Vercel Deployment
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE_PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "05  Vercel + Supabase Production Deployment", font_size=28, color=BLUE_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05), BLUE_PRIMARY)

# Steps
vercel_steps = [
    ("1", "Push to GitHub", "Push your code to a GitHub/GitLab/Bitbucket repository", BLUE_PRIMARY),
    ("2", "Import to Vercel", "vercel.com > New Project > Import Git Repository", TEAL),
    ("3", "Configure Settings", "Framework: Next.js (auto-detected)\nRoot Directory: ./ (default)\nBuild Command: next build\nOutput Directory: .next", GREEN),
    ("4", "Set Environment Variables", "Add ALL 4 env vars in Vercel Dashboard:\nSUPABASE_URL, SUPABASE_ANON_KEY,\nSUPABASE_SERVICE_ROLE_KEY,\nNEXT_PUBLIC_BASE_URL=https://your-domain.vercel.app", ORANGE),
    ("5", "Deploy!", "Click Deploy > Wait for build > Verify at your URL\nFirst visit: data auto-seeds via /api/seed", BLUE_PRIMARY),
]

for i, (num, title, desc, color) in enumerate(vercel_steps):
    y = Inches(1.5) + Inches(i * 1.1)
    add_shape(slide, Inches(0.8), y, Inches(0.5), Inches(0.5), color)
    add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(0.5), Inches(0.4), num, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), y + Inches(0.05), Inches(4), Inches(0.4), title, font_size=16, color=GRAY_900, bold=True)
    add_text_box(slide, Inches(1.5), y + Inches(0.35), Inches(5.5), Inches(0.7), desc, font_size=11, color=GRAY_500)

# Right side - Vercel config
add_shape(slide, Inches(7.5), Inches(1.5), Inches(5.1), Inches(5.5), GRAY_100)
add_text_box(slide, Inches(7.8), Inches(1.7), Inches(4.5), Inches(0.4),
             "Vercel Environment Variables", font_size=16, color=GRAY_900, bold=True)
add_code_block(slide, Inches(7.8), Inches(2.2), Inches(4.5), Inches(2.5),
"""# In Vercel Dashboard > Settings > Environment Variables

SUPABASE_URL=https://xxx.supabase.co
SUPABASE_ANON_KEY=eyJhbGciOi...
SUPABASE_SERVICE_ROLE_KEY=eyJhbGciOi...
NEXT_PUBLIC_BASE_URL=https://your-app.vercel.app

# Important: Set for ALL environments
# (Production, Preview, Development)""", font_size=10)

add_text_box(slide, Inches(7.8), Inches(4.9), Inches(4.5), Inches(0.4),
             "Custom Domain (Optional)", font_size=14, color=GRAY_900, bold=True)
add_text_box(slide, Inches(7.8), Inches(5.3), Inches(4.5), Inches(1.5),
             "1. Vercel > Settings > Domains\n2. Add your domain\n3. Update DNS records (CNAME/A)\n4. SSL auto-provisioned by Vercel\n5. Update NEXT_PUBLIC_BASE_URL",
             font_size=11, color=GRAY_500)

# ============================================================
# SLIDE 8: Environment Variables Reference
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE_PRIMARY)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "06  Environment Variables Reference", font_size=28, color=BLUE_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05), BLUE_PRIMARY)

# Table header
add_shape(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5), BLUE_PRIMARY)
headers = [("Variable", 3.2), ("Example Value", 4.5), ("Required", 1.0), ("Where to Find", 3.0)]
x = Inches(0.8)
for header, width in headers:
    add_text_box(slide, x + Inches(0.1), Inches(1.55), Inches(width), Inches(0.4),
                 header, font_size=12, color=WHITE, bold=True)
    x += Inches(width)

# Table rows
env_vars = [
    ("SUPABASE_URL", "https://xxx.supabase.co", "Yes", "Supabase > Settings > API"),
    ("SUPABASE_ANON_KEY", "eyJhbG...(public JWT)", "Yes", "Supabase > Settings > API"),
    ("SUPABASE_SERVICE_ROLE_KEY", "eyJhbG...(secret JWT)", "Yes", "Supabase > Settings > API"),
    ("NEXT_PUBLIC_BASE_URL", "https://your-domain.com", "Yes", "Your deployment URL"),
    ("CORS_ORIGINS", "*", "No", "Set to domain in production"),
]

for i, (var, example, required, where) in enumerate(env_vars):
    y = Inches(2.05) + Inches(i * 0.55)
    bg = GRAY_100 if i % 2 == 0 else WHITE
    add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.5), bg)
    x = Inches(0.8)
    vals = [(var, 3.2, GRAY_900, True), (example, 4.5, GRAY_500, False), (required, 1.0, GREEN if required == "Yes" else GRAY_500, True), (where, 3.0, GRAY_500, False)]
    for val, width, color, bold in vals:
        add_text_box(slide, x + Inches(0.1), y + Inches(0.08), Inches(width), Inches(0.35),
                     val, font_size=11, color=color, bold=bold)
        x += Inches(width)

# Security note
add_shape(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.8), RGBColor(0xFE, 0xE2, 0xE2))
add_text_box(slide, Inches(1.1), Inches(5.2), Inches(11), Inches(0.4),
             "SECURITY WARNINGS", font_size=16, color=RED, bold=True)
warnings = [
    "NEVER commit .env files to version control - add .env to .gitignore",
    "NEVER expose SUPABASE_SERVICE_ROLE_KEY in client-side code (no NEXT_PUBLIC_ prefix!)",
    "The service_role key bypasses Row Level Security - keep it server-side only",
    "In production, set CORS_ORIGINS to your specific domain instead of '*'",
    "Rotate keys immediately if accidentally exposed in public repositories",
]
add_bullet_list(slide, Inches(1.1), Inches(5.6), Inches(11), Inches(1.2),
                [f"  {w}" for w in warnings], font_size=11, color=RED, spacing=Pt(3))

# ============================================================
# SLIDE 9: Common Errors (Part 1)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), RED)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "07  Common Errors & Fixes (Part 1)", font_size=28, color=BLUE_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05), RED)

errors_1 = [
    ("relation \"departments\" does not exist",
     'PostgrestError: relation "public.departments" does not exist',
     "Tables not created yet. Run the SQL script in Supabase SQL Editor (Dashboard > SQL Editor > New Query > Paste SQL > Run)"),
    ("SUPABASE_URL is not defined",
     "TypeError: Cannot read properties of undefined (reading 'from')\ncreateClient requires a valid URL",
     "Missing .env file or env vars. Create .env with SUPABASE_URL, SUPABASE_ANON_KEY, SUPABASE_SERVICE_ROLE_KEY. Restart server after changes."),
    ("JWT expired / Invalid API key",
     '{"message":"Invalid API key","hint":"Double check your Supabase URL and key"}',
     "Verify your Supabase keys are correct. Check Settings > API in Supabase Dashboard. Ensure no extra spaces/newlines in .env values."),
]

y_offset = Inches(1.5)
for i, (title, msg, solution) in enumerate(errors_1):
    h = add_error_card(slide, Inches(0.8), y_offset, Inches(11.7), title, msg, solution, i)
    y_offset += h + Inches(0.15)

# ============================================================
# SLIDE 10: Common Errors (Part 2)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), RED)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "07  Common Errors & Fixes (Part 2)", font_size=28, color=BLUE_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05), RED)

errors_2 = [
    ("Vercel Build Fails - Module Not Found",
     "Module not found: Can't resolve '@supabase/supabase-js'\nError: Build failed with exit code 1",
     "Ensure @supabase/supabase-js is in package.json dependencies. Run 'yarn install' locally first, commit yarn.lock file, then redeploy."),
    ("CORS / Network Error on API Calls",
     "Access-Control-Allow-Origin header missing\nFailed to fetch /api/departments",
     "Check NEXT_PUBLIC_BASE_URL matches your actual domain. API routes must use /api prefix. In Vercel, redeploy after adding env vars."),
    ("duplicate key value violates unique constraint",
     'error: duplicate key value violates unique constraint "users_email_key"',
     "Database already has this record. The /api/seed endpoint checks for existing data. Clear the table in Supabase Dashboard > Table Editor, or skip seeding."),
    ("Vercel Serverless Function Timeout",
     "FUNCTION_INVOCATION_TIMEOUT: Task timed out after 10 seconds",
     "Supabase free tier cold starts can be slow. Upgrade to Pro plan, or add connection retry logic. Check Supabase project isn't paused (free tier pauses after 7 days of inactivity)."),
]

y_offset = Inches(1.5)
for i, (title, msg, solution) in enumerate(errors_2):
    h = add_error_card(slide, Inches(0.8), y_offset, Inches(11.7), title, msg, solution, i + 3)
    y_offset += h + Inches(0.1)

# ============================================================
# SLIDE 11: Common Errors (Part 3)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), RED)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "07  Common Errors & Fixes (Part 3)", font_size=28, color=BLUE_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05), RED)

errors_3 = [
    ("Hydration Mismatch (Next.js)",
     "Error: Hydration failed because server rendered HTML didn't match client.\nText content does not match.",
     "This happens with localStorage access during SSR. Use 'use client' directive. Check typeof window !== 'undefined' before localStorage. Already handled in this project."),
    ("RLS (Row Level Security) Blocking Queries",
     '{"code":"42501","message":"permission denied for table departments"}',
     "Using service_role key bypasses RLS. If using anon key, add RLS policies: ALTER TABLE departments ENABLE ROW LEVEL SECURITY; CREATE POLICY ... Using service role key is recommended for API routes."),
    ("Supabase Project Paused (Free Tier)",
     "FetchError: request to https://xxx.supabase.co/rest/v1/departments failed\nconnection refused / timeout",
     "Free Supabase projects pause after 7 days of inactivity. Go to Supabase Dashboard > Resume project. Consider upgrading to Pro ($25/mo) for always-on."),
    ("next build fails: ESLint errors",
     "Failed to compile. ESLint: 'X' is defined but never used.\nnext build failed with exit code 1",
     "Add to next.config.js: eslint: { ignoreDuringBuilds: true }. Or fix the specific ESLint errors shown in the build output."),
]

y_offset = Inches(1.5)
for i, (title, msg, solution) in enumerate(errors_3):
    h = add_error_card(slide, Inches(0.8), y_offset, Inches(11.7), title, msg, solution, i + 7)
    y_offset += h + Inches(0.1)

# ============================================================
# SLIDE 12: Post-Deployment Checklist
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GREEN)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "08  Post-Deployment Checklist", font_size=28, color=BLUE_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(2), Inches(0.05), GREEN)

checklist_left = [
    ("Functionality Checks", [
        "Home page loads with departments & doctors",
        "Appointment booking flow works end-to-end",
        "Contact form submits successfully",
        "Admin login works (admin@medcare.com / admin123)",
        "Admin can manage appointments (complete/cancel/delete)",
        "Admin can view & delete leads with phone numbers",
        "Blog articles load and display correctly",
        "All navigation links work properly",
    ]),
    ("Performance", [
        "Page load time < 3 seconds",
        "API response time < 500ms",
        "Images lazy-loaded properly",
        "No console errors in browser",
    ]),
]

checklist_right = [
    ("Security", [
        "HTTPS enabled (auto on Vercel)",
        ".env not committed to git",
        "Service role key only on server-side",
        "Input validation on all forms",
        "CORS configured for production domain",
    ]),
    ("SEO & Monitoring", [
        "Page title & meta description set",
        "Open Graph tags configured",
        "Verify Google Search Console",
        "Set up Vercel Analytics",
        "Monitor Supabase usage dashboard",
        "Set up uptime monitoring (optional)",
    ]),
]

for col_data, x_start in [(checklist_left, Inches(0.8)), (checklist_right, Inches(7.0))]:
    y = Inches(1.5)
    for section_title, items in col_data:
        add_text_box(slide, x_start, y, Inches(5.5), Inches(0.4),
                     section_title, font_size=16, color=GRAY_900, bold=True)
        y += Inches(0.4)
        for item in items:
            add_text_box(slide, x_start + Inches(0.1), y, Inches(5.3), Inches(0.3),
                         f"   {item}", font_size=11, color=GRAY_700)
            y += Inches(0.28)
        y += Inches(0.2)

# ============================================================
# SLIDE 13: Quick Reference / Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.15), BLUE_PRIMARY)

add_text_box(slide, Inches(1.5), Inches(0.8), Inches(10.3), Inches(0.8),
             "Quick Reference Card", font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Quick commands
add_shape(slide, Inches(0.8), Inches(1.8), Inches(3.7), Inches(4.5), RGBColor(0x1D, 0x2D, 0x50))
add_text_box(slide, Inches(1.0), Inches(2.0), Inches(3.3), Inches(0.4),
             "Local Commands", font_size=16, color=BLUE_ACCENT, bold=True)
add_code_block(slide, Inches(1.0), Inches(2.5), Inches(3.3), Inches(3.5),
"""# Install
yarn install

# Development
yarn dev

# Production build
yarn build
yarn start

# Lint
yarn lint

# Seed database
curl localhost:3000/api/seed""", font_size=10)

# Vercel commands
add_shape(slide, Inches(4.8), Inches(1.8), Inches(3.7), Inches(4.5), RGBColor(0x1D, 0x2D, 0x50))
add_text_box(slide, Inches(5.0), Inches(2.0), Inches(3.3), Inches(0.4),
             "Vercel Commands", font_size=16, color=BLUE_ACCENT, bold=True)
add_code_block(slide, Inches(5.0), Inches(2.5), Inches(3.3), Inches(3.5),
"""# Install Vercel CLI
npm i -g vercel

# Deploy from CLI
vercel

# Production deploy
vercel --prod

# Check logs
vercel logs

# Environment vars
vercel env add
vercel env ls""", font_size=10)

# Key URLs
add_shape(slide, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.5), RGBColor(0x1D, 0x2D, 0x50))
add_text_box(slide, Inches(9.0), Inches(2.0), Inches(3.3), Inches(0.4),
             "Key URLs", font_size=16, color=BLUE_ACCENT, bold=True)
urls = [
    "API Health Check:",
    "  /api/health",
    "",
    "Seed Database:",
    "  /api/seed",
    "",
    "Get Setup SQL:",
    "  /api/setup-sql",
    "",
    "Admin Login:",
    "  admin@medcare.com",
    "  admin123",
    "",
    "Supabase Dashboard:",
    "  supabase.com/dashboard",
]
add_bullet_list(slide, Inches(9.0), Inches(2.5), Inches(3.3), Inches(3.5),
                urls, font_size=10, color=RGBColor(0xA6, 0xE2, 0x2E), spacing=Pt(2))

# Footer text
add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.5),
             "MedCare Multi-Speciality Hospital  |  Built with Next.js + Supabase  |  2025",
             font_size=12, color=GRAY_500, alignment=PP_ALIGN.CENTER)

# ============== SAVE ==============
output_path = "/app/public/MedCare_Deployment_Guide.pptx"
prs.save(output_path)
print(f"PPT saved to {output_path}")
print(f"Download URL: /MedCare_Deployment_Guide.pptx")
